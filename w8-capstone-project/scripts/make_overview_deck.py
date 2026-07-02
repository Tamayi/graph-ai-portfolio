"""Generate the capstone overview deck (docs/SitRep_Intelligence_Overview_v3.pptx).

Rebuilds the v2 deck content with a cleaner design system, plus two new
slides: "How AI is used" (development + product features) and a GenAI
concepts map for a non-developer audience. Run:

    python scripts/make_overview_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "SitRep_Intelligence_Overview_v3.pptx"

# --- Design system -----------------------------------------------------------
W, H = Inches(13.333), Inches(7.5)

PAPER = RGBColor.from_string("FBFAF6")   # warm off-white background
INK = RGBColor.from_string("1E2B2F")     # near-black text
MUTED = RGBColor.from_string("5E6E73")   # secondary text
TEAL = RGBColor.from_string("0B4F5C")    # primary brand
TEAL_DARK = RGBColor.from_string("07333C")
AMBER = RGBColor.from_string("D97E00")   # accent
CARD = RGBColor.from_string("FFFFFF")
BORDER = RGBColor.from_string("E4DFD3")
TEAL_TINT = RGBColor.from_string("E7F0F1")

DISPLAY = "Georgia"      # headings
BODY = "Segoe UI"        # everything else

FOOTER = "MVE / Ebola SitRep Intelligence  ·  INSP DRC"


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def paint_bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill=CARD, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.06):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=None):
    """Add a textbox. `runs` is a list of paragraphs; each paragraph is a list
    of (text, size, bold, color, font, italic) tuples (trailing items optional).
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = space_after
        for item in para:
            t, size = item[0], item[1]
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else INK
            font = item[4] if len(item) > 4 else BODY
            italic = item[5] if len(item) > 5 else False
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
            r.font.italic = italic
    return tb


def footer(slide, n):
    text(slide, Inches(0.55), Inches(7.08), Inches(8), Inches(0.35),
         [[(FOOTER, 9, False, MUTED)]])
    text(slide, Inches(12.35), Inches(7.08), Inches(0.5), Inches(0.35),
         [[(str(n), 9, True, MUTED)]], align=PP_ALIGN.RIGHT)


def header(slide, number, kicker, title, subtitle=None):
    """Section header: amber number chip, kicker, big serif title, rule."""
    chip = box(slide, Inches(0.55), Inches(0.5), Inches(0.62), Inches(0.62),
               fill=TEAL, line=None, radius=0.18)
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = number
    r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("FFFFFF"); r.font.name = DISPLAY
    text(slide, Inches(1.4), Inches(0.47), Inches(11.3), Inches(0.35),
         [[(kicker.upper(), 12, True, AMBER)]])
    text(slide, Inches(1.4), Inches(0.78), Inches(11.3), Inches(0.6),
         [[(title, 26, True, INK, DISPLAY)]])
    y = 1.42
    if subtitle:
        text(slide, Inches(1.4), Inches(1.38), Inches(11.3), Inches(0.4),
             [[(subtitle, 13, False, MUTED)]])
        y = 1.82
    ln = box(slide, Inches(0.55), Inches(y), Inches(12.23), Pt(1.6),
             fill=BORDER, line=None, shape=MSO_SHAPE.RECTANGLE)
    return Inches(y + 0.22)


def card(slide, x, y, w, h, title, body, title_color=TEAL, accent=True,
         title_size=13.5, body_size=11):
    box(slide, x, y, w, h)
    if accent:
        box(slide, x, y, Inches(0.06), h, fill=AMBER, line=None,
            shape=MSO_SHAPE.RECTANGLE)
    pad = Inches(0.18)
    text(slide, x + pad + Inches(0.04), y + Inches(0.14), w - pad * 2, Inches(0.4),
         [[(title, title_size, True, title_color)]])
    text(slide, x + pad + Inches(0.04), y + Inches(0.52), w - pad * 2,
         h - Inches(0.66), [[(body, body_size, False, INK)]])


def arrow_right(slide, x, y, w=Inches(0.34), h=Inches(0.22)):
    box(slide, x, y, w, h, fill=AMBER, line=None, shape=MSO_SHAPE.RIGHT_ARROW)


def flow_box(slide, x, y, w, h, title, sub, fill=CARD, tcolor=TEAL,
             scolor=MUTED):
    box(slide, x, y, w, h, fill=fill,
        line=None if fill != CARD else BORDER)
    text(slide, x + Inches(0.12), y, w - Inches(0.24), h,
         [[(title, 12, True, tcolor)], [(sub, 9.5, False, scolor)]],
         anchor=MSO_ANCHOR.MIDDLE)


# ==============================================================================
prs = Presentation()
prs.slide_width, prs.slide_height = W, H

# --- 1 · Title ---------------------------------------------------------------
s = blank_slide(prs)
paint_bg(s, TEAL_DARK)
box(s, 0, 0, W, Inches(0.14), fill=AMBER, line=None, shape=MSO_SHAPE.RECTANGLE)
text(s, Inches(1.0), Inches(1.15), Inches(11.3), Inches(0.4),
     [[("OUTBREAK SITUATION-REPORT INTELLIGENCE", 13, True, AMBER)]])
text(s, Inches(1.0), Inches(1.7), Inches(11.3), Inches(2.1),
     [[("From PDF SitReps to", 44, True, RGBColor.from_string("FFFFFF"), DISPLAY)],
      [("Verified Outbreak Intelligence", 44, True, AMBER, DISPLAY)]])
text(s, Inches(1.0), Inches(3.75), Inches(10.5), Inches(0.8),
     [[("A grounded RAG system over the 17th Ebola (MVE/BVD) outbreak reports "
        "from the Institut National de Santé Publique, DR Congo.",
        16, False, RGBColor.from_string("D9E4E6"))]])
chips = ["Ask & cite", "Verify quality", "Track evolution", "Surface gaps"]
cx = 1.0
for c in chips:
    wch = 0.32 + 0.105 * len(c)
    ch = box(s, Inches(cx), Inches(4.75), Inches(wch), Inches(0.5),
             fill=None, line=AMBER, radius=0.5)
    tf = ch.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = c
    r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("FFFFFF"); r.font.name = BODY
    cx += wch + 0.25
text(s, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.4),
     [[("Claude  ·  Voyage embeddings  ·  Chroma  ·  FastAPI  ·  Python 3.13",
        12.5, False, RGBColor.from_string("9FB8BD"))]])

# --- 2 · The problem -----------------------------------------------------------
s = blank_slide(prs); paint_bg(s)
top = header(s, "01", "The problem",
             "The pulse of an outbreak — but hard to use",
             "INSP publishes a new PDF almost every day. Analysts must read fast, "
             "trust the numbers, and see the trend — manually, under pressure.")
cards = [
    ("Data quality", "Typos, arithmetic that doesn't add up, missing tables, "
     "inconsistent CFR%."),
    ("Verification", "Are the figures internally consistent? Do today's totals "
     "agree with yesterday's?"),
    ("Evolution over time", "40+ reports, N°001 → N°042. How has the outbreak "
     "actually moved?"),
    ("Interaction", "No way to just ask a question and get an answer with a "
     "page-level citation."),
]
cw, gap = Inches(2.92), Inches(0.18)
for i, (t_, b_) in enumerate(cards):
    card(s, Inches(0.55) + i * (cw + gap), top + Inches(0.15), cw, Inches(2.5),
         t_, b_)
band = box(s, Inches(0.55), top + Inches(3.0), Inches(12.23), Inches(1.15),
           fill=TEAL_TINT, line=None)
text(s, Inches(0.85), top + Inches(3.16), Inches(11.6), Inches(0.9),
     [[("Consequence  ", 13, True, TEAL),
       ("Decisions ride on numbers nobody has had time to validate, and trends "
        "live only in an analyst's memory of yesterday's PDF.", 13, False, INK)]])
footer(s, 2)

# --- 3 · The solution ----------------------------------------------------------
s = blank_slide(prs); paint_bg(s)
top = header(s, "02", "The solution",
             "One workspace: read every report, check it, ask it")
cards = [
    ("Ground every answer", "Responses come only from retrieved report "
     "excerpts, each cited as SitRep N° · date · page."),
    ("Trust the numbers", "Automatic completeness, consistency and arithmetic "
     "checks on each report — and between reports."),
    ("See the trajectory", "Structured 'facts' datasets turn 40 PDFs into "
     "clean KPI / province / zone time-series."),
    ("Stay current", "Drop in a new INSP link or PDF and it is parsed, "
     "extracted, indexed and ready to query."),
]
cw = Inches(2.92)
for i, (t_, b_) in enumerate(cards):
    card(s, Inches(0.55) + i * (cw + gap), top + Inches(0.15), cw, Inches(2.35), t_, b_)
box(s, Inches(0.55), top + Inches(2.8), Inches(12.23), Inches(2.0),
    fill=TEAL, line=None)
text(s, Inches(0.9), top + Inches(2.98), Inches(11.6), Inches(0.35),
     [[("HOW IT STAYS HONEST", 12, True, AMBER)]])
text(s, Inches(0.9), top + Inches(3.35), Inches(11.6), Inches(1.4),
     [[("Structured-first", 11.5, True, RGBColor.from_string("FFFFFF")),
       ("  exact figures come from JSON datasets, not the model's guess.      ",
        11.5, False, RGBColor.from_string("D9E4E6")),
       ("Nullable by design", 11.5, True, RGBColor.from_string("FFFFFF")),
       ("  a missing table stays null; the system never invents a number.",
        11.5, False, RGBColor.from_string("D9E4E6"))],
      [("Tool-forced output", 11.5, True, RGBColor.from_string("FFFFFF")),
       ("  assessments & comparisons return as validated JSON, not prose.      ",
        11.5, False, RGBColor.from_string("D9E4E6")),
       ("Holdout set", 11.5, True, RGBColor.from_string("FFFFFF")),
       ("  recent reports can be withheld from the index to test retrieval honestly.",
        11.5, False, RGBColor.from_string("D9E4E6"))]],
     space_after=Pt(8))
footer(s, 3)

# --- 4 · How AI is used (NEW) ---------------------------------------------------
s = blank_slide(prs); paint_bg(s)
top = header(s, "03", "How AI is used",
             "AI built the application — and AI runs inside it",
             "Two distinct roles: a development partner during the build, and "
             "generative features the analyst uses every day.")
colw = Inches(6.0)
# Left column: building with AI
box(s, Inches(0.55), top + Inches(0.1), colw, Inches(4.5), fill=CARD, line=BORDER)
text(s, Inches(0.85), top + Inches(0.3), colw - Inches(0.6), Inches(0.4),
     [[("BUILDING THE APP  ·  AI AS DEVELOPMENT PARTNER", 12, True, AMBER)]])
left_items = [
    ("Pair-programming with Claude Code",
     "Backend, ingestion pipeline and tests scaffolded and iterated "
     "conversationally, with the developer reviewing every change."),
    ("AI-designed frontend & translations",
     "The dashboard screens were generated with Claude Design, and the "
     "FR / EN / PT interface translations were AI-generated too."),
    ("Prompt engineering as code",
     "All LLM instructions live as versioned templates in /prompts — reviewed, "
     "diffed and improved like any other source file."),
    ("Tests without the network",
     "AI-authored pytest suite with mocked API calls, so quality checks run "
     "free and fast on every change."),
]
yy = top + Inches(0.72)
for t_, b_ in left_items:
    text(s, Inches(0.85), yy, colw - Inches(0.6), Inches(0.95),
         [[("▸ ", 11.5, True, AMBER), (t_, 11.5, True, INK)],
          [(b_, 10.5, False, MUTED)]])
    yy += Inches(0.94)
# Right column: GenAI in the product
box(s, Inches(6.78), top + Inches(0.1), colw, Inches(4.5), fill=TEAL, line=None)
text(s, Inches(7.08), top + Inches(0.3), colw - Inches(0.6), Inches(0.4),
     [[("INSIDE THE PRODUCT  ·  GENERATIVE AI FEATURES", 12, True, AMBER)]])
right_items = [
    ("Document understanding", "Claude reads each PDF report and extracts a "
     "structured, nullable 'facts' record."),
    ("Semantic search (RAG)", "Voyage AI embeddings + Chroma retrieve the "
     "most relevant excerpts for any question."),
    ("Grounded, cited chat", "Multi-turn Q&A in FR/EN/PT; every figure cites "
     "SitRep N°, date and page."),
    ("Function calling", "The model calls query_data for exact numbers and "
     "plot_chart to draw trends."),
    ("LLM as analyst", "Quality scoring, two-report comparison and topic "
     "synthesis are AI-generated, human-verified."),
]
yy = top + Inches(0.72)
for t_, b_ in right_items:
    text(s, Inches(7.08), yy, colw - Inches(0.6), Inches(0.8),
         [[("▸ ", 11.5, True, AMBER),
           (t_, 11.5, True, RGBColor.from_string("FFFFFF"))],
          [(b_, 10.5, False, RGBColor.from_string("C8D9DC"))]])
    yy += Inches(0.75)
footer(s, 4)

# --- 5 · GenAI concepts map (NEW) ------------------------------------------------
s = blank_slide(prs); paint_bg(s)
top = header(s, "04", "GenAI concepts, applied",
             "Every core concept from the course, at work in one system")
concepts = [
    ("RAG", "Retrieve the 8 most relevant excerpts, answer only from them — "
     "the model cannot drift from the source."),
    ("Embeddings", "voyage-4 multilingual vectors put French reports and "
     "English questions in the same semantic space."),
    ("Grounding & citations", "Every answer carries SitRep N° · date · page, "
     "so a human can check it in seconds."),
    ("Tool use / function calling", "Claude decides when to call query_data "
     "(exact figures) or plot_chart (visuals)."),
    ("Structured output", "Extraction and assessments are schema-validated "
     "JSON — nullable, never invented."),
    ("Prompt engineering", "Externalised, versioned prompt templates; "
     "per-language rendering (FR/EN/PT)."),
    ("Evaluation & holdout", "Recent reports can be withheld from the index "
     "to test retrieval honestly."),
    ("Token efficiency", "Prompt caching plus history compaction keep "
     "multi-turn chat fast and affordable."),
]
cw, chh = Inches(2.92), Inches(1.95)
for i, (t_, b_) in enumerate(concepts):
    row, col = divmod(i, 4)
    card(s, Inches(0.55) + col * (cw + gap),
         top + Inches(0.15) + row * (chh + Inches(0.22)),
         cw, chh, t_, b_, title_size=12.5, body_size=10.5)
footer(s, 5)

# --- 6 · System components -------------------------------------------------------
s = blank_slide(prs); paint_bg(s)
top = header(s, "05", "System components",
             "Six capabilities over one grounded data core")
comp = [
    ("Grounded chat assistant", "Multi-turn Q&A in FR / EN / PT. Tools: "
     "query_data (exact figures) + plot_chart (trends). Always cited."),
    ("Outbreak summary + epi analysis", "KPI, per-province & per-zone "
     "time-series: cases, deaths, CFR%, labs, contacts, IPC, vaccination."),
    ("Single-document reviewer", "Scores 1 report on 5 dimensions, flags "
     "arithmetic & consistency issues by severity."),
    ("Document comparison", "Validates two reports: totals must not drop; "
     "flags suspect jumps, contradictions, section changes."),
    ("Topics analysis", "Extracts gaps & challenges per report, themes them, "
     "synthesises recurring bottlenecks."),
    ("SitRep library", "Add a report via INSP link, PDF URL or upload → auto "
     "parse, extract, index, aggregate."),
]
cw, chh = Inches(3.94), Inches(1.72)
for i, (t_, b_) in enumerate(comp):
    row, col = divmod(i, 3)
    card(s, Inches(0.55) + col * (cw + Inches(0.2)),
         top + Inches(0.1) + row * (chh + Inches(0.2)),
         cw, chh, t_, b_, title_size=12.5, body_size=10.5)
box(s, Inches(0.55), top + Inches(4.0), Inches(12.23), Inches(0.95),
    fill=TEAL_TINT, line=None)
text(s, Inches(0.85), top + Inches(4.12), Inches(11.7), Inches(0.75),
     [[("The data core beneath them all  ", 12, True, TEAL),
       ("JSON 'facts' datasets (exact numbers)  +  Chroma vector store "
        "(semantic search)  +  externalised prompts  +  Claude for extraction, "
        "analysis & synthesis.", 12, False, INK)]])
footer(s, 6)

# --- 7 · Architecture & data flow ---------------------------------------------
s = blank_slide(prs); paint_bg(s)
top = header(s, "06", "Architecture & data flow",
             "Ingest once, serve many",
             "Structured facts and semantic text, side by side — built once "
             "at ingest, queried everywhere.")
fy = top + Inches(0.78)   # centre row; leaves room for the raised store box
fh = Inches(0.95)
flow_box(s, Inches(0.55), fy, Inches(2.15), fh, "Source SitReps",
         "INSP PDFs / JSON bundle")
arrow_right(s, Inches(2.78), fy + Inches(0.36))
flow_box(s, Inches(3.2), fy, Inches(2.15), fh, "Ingest pipeline",
         "parse · chunk · page markers")
arrow_right(s, Inches(5.43), fy + Inches(0.36))
# two parallel stores
flow_box(s, Inches(5.85), fy - Inches(0.62), Inches(2.7), fh,
         "Structured facts (JSON)", "KPI · province · zone · pillars",
         fill=TEAL_TINT)
flow_box(s, Inches(5.85), fy + Inches(0.62), Inches(2.7), fh,
         "Vector store (Chroma)", "Voyage embeddings · top-k",
         fill=TEAL_TINT)
arrow_right(s, Inches(8.63), fy + Inches(0.36))
flow_box(s, Inches(9.05), fy, Inches(3.7), fh, "Claude",
         "extract · assess · compare · synthesise", fill=TEAL,
         tcolor=RGBColor.from_string("FFFFFF"),
         scolor=RGBColor.from_string("C8D9DC"))
# serving row
sy = fy + Inches(1.75)
box(s, Inches(0.55), sy, Inches(12.2), Inches(0.85), fill=CARD, line=BORDER)
text(s, Inches(0.85), sy, Inches(3.0), Inches(0.85),
     [[("FastAPI backend · /api/*", 12.5, True, TEAL)]],
     anchor=MSO_ANCHOR.MIDDLE)
eps = ["/chat", "/overview", "/data-quality/assess", "/data-quality/compare",
       "/topics/run", "/sitreps"]
ex = 3.7
for ep in eps:
    wch = 0.4 + 0.075 * len(ep)
    c = box(s, Inches(ex), sy + Inches(0.2), Inches(wch), Inches(0.45),
            fill=TEAL_TINT, line=None, radius=0.5)
    tf = c.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ep
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = TEAL
    r.font.name = "Consolas"
    ex += wch + 0.12
dy = sy + Inches(0.97)
box(s, Inches(0.55), dy, Inches(12.2), Inches(0.6), fill=CARD, line=BORDER)
text(s, Inches(0.85), dy, Inches(11.6), Inches(0.6),
     [[("Frontend  ", 12.5, True, TEAL),
       ("6 screens · FR/EN/PT · live API data · Plotly charts",
        11.5, False, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(0.55), dy + Inches(0.72), Inches(12.2), Inches(0.55),
    fill=TEAL_TINT, line=None)
text(s, Inches(0.85), dy + Inches(0.72), Inches(11.6), Inches(0.55),
     [[("Caching & incrementality  ", 11.5, True, TEAL),
       ("per-report extraction is cached — a new upload only extracts itself, "
        "then re-aggregates the datasets.", 11.5, False, INK)]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 7)

# --- 8 · Topics analysis ---------------------------------------------------------
s = blank_slide(prs); paint_bg(s)
top = header(s, "07", "Topics analysis",
             "Turning narrative 'défis & lacunes' into a strategic picture",
             "Each report buries its real story in prose — the gaps, constraints "
             "and unmet needs. The system extracts, themes, and synthesises them.")
card(s, Inches(0.55), top + Inches(0.15), Inches(5.9), Inches(1.75),
     "1 · Extract  (gaps_extract.md)",
     "Claude reads each SitRep and returns JSON items:\n"
     "{ sitrep_number, date, theme, detail }", body_size=11.5)
card(s, Inches(0.55), top + Inches(2.1), Inches(5.9), Inches(1.75),
     "2 · Synthesise  (gaps_synthesise.md)",
     "Aggregates all items into a narrative of the chronic bottlenecks, "
     "cached for reuse.", body_size=11.5)
box(s, Inches(6.75), top + Inches(0.15), Inches(6.03), Inches(3.7),
    fill=TEAL, line=None)
text(s, Inches(7.05), top + Inches(0.35), Inches(5.4), Inches(0.35),
     [[("THEMES TRACKED", 12, True, AMBER)]])
themes = ["Financement", "Surveillance", "Vaccination", "Sécurité",
          "Logistique", "Engagement com.", "Ressources hum.", "Laboratoire"]
tx, ty = 7.05, None
row_y = top + Inches(0.85)
for i, th in enumerate(themes):
    row, col = divmod(i, 2)
    c = box(s, Inches(7.05) + col * Inches(2.85),
            row_y + row * Inches(0.62), Inches(2.65), Inches(0.48),
            fill=None, line=RGBColor.from_string("6FA7B0"), radius=0.5)
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = th
    r.font.size = Pt(11.5); r.font.color.rgb = RGBColor.from_string("FFFFFF")
    r.font.name = BODY
box(s, Inches(0.55), top + Inches(4.05), Inches(12.23), Inches(0.85),
    fill=TEAL_TINT, line=None)
text(s, Inches(0.85), top + Inches(4.17), Inches(11.6), Inches(0.6),
     [[("Payoff  ", 12, True, TEAL),
       ("a coordinator sees that financing and logistics gaps recur for weeks — "
        "evidence for where to push, not a vague impression.", 12, False, INK)]])
footer(s, 8)

# --- 9 · Tech stack --------------------------------------------------------------
s = blank_slide(prs); paint_bg(s)
top = header(s, "08", "Tech stack", "Lean, mostly-local, API-light")
stack = [
    ("AI & language", "Anthropic Claude — Sonnet default,\nOpus for the "
     "hardest tasks\nTool-forcing for structured JSON\nPrompts externalised "
     "in /prompts"),
    ("Retrieval", "Voyage AI multilingual embeddings\nChroma local persistent "
     "vector DB\n~1400-char chunks, 200 overlap\nTop-k = 8, page-level "
     "citations"),
    ("Backend & data", "FastAPI + Uvicorn (async)\npypdf text extraction\n"
     "pandas · Plotly\nJSON datasets + manifest CSV"),
    ("Platform & quality", "Python 3.13, uv package manager\npytest "
     "(network-free mocks)\nruff + black, Makefile tasks\nClaude-Design HTML "
     "frontend"),
]
cw = Inches(2.92)
for i, (t_, b_) in enumerate(stack):
    card(s, Inches(0.55) + i * (cw + gap), top + Inches(0.2), cw, Inches(3.3),
         t_, b_, body_size=11)
footer(s, 9)

# --- 10 · Limitations & roadmap ---------------------------------------------------
s = blank_slide(prs); paint_bg(s)
top = header(s, "09", "Limitations & roadmap",
             "Honest about today's edges — and what's next")
box(s, Inches(0.55), top + Inches(0.15), Inches(6.0), Inches(4.55),
    fill=CARD, line=BORDER)
text(s, Inches(0.85), top + Inches(0.33), Inches(5.4), Inches(0.4),
     [[("CURRENT LIMITATIONS", 12, True, MUTED)]])
lims = [
    ("API keys & cost", "needs Anthropic + Voyage keys; LLM and embedding "
     "calls are metered usage."),
    ("No OCR", "pypdf reads text only — scanned or image-only PDFs won't parse."),
    ("Format-coupled", "facts schema & prompts are tuned to the INSP / MVE "
     "report layout."),
    ("Human-in-the-loop", "LLM extraction & scoring assist analysts — they "
     "aren't authoritative."),
    ("Single-user & local", "local Chroma store; no auth, access control or "
     "multi-tenant serving."),
]
yy = top + Inches(0.75)
for t_, b_ in lims:
    text(s, Inches(0.85), yy, Inches(5.45), Inches(0.75),
         [[(t_ + "  ", 11.5, True, INK), (b_, 11, False, MUTED)]])
    yy += Inches(0.74)
box(s, Inches(6.78), top + Inches(0.15), Inches(6.0), Inches(4.55),
    fill=TEAL, line=None)
text(s, Inches(7.08), top + Inches(0.33), Inches(5.4), Inches(0.4),
     [[("FUTURE DEVELOPMENT", 12, True, AMBER)]])
future = [
    ("Multi-outbreak support", "generalise schema & prompts so any outbreak "
     "or agency can plug in."),
    ("One-click DOCX export", "share a quality assessment or comparison as "
     "Word — no copy-paste."),
    ("Deeper text analysis", "intervention timelines, partner engagement, "
     "response over time."),
    ("Bring-your-own key", "users supply their own API keys and pick models "
     "per deployment."),
    ("Hardening", "auth, OCR fallback, and formal retrieval benchmarking on "
     "the holdout set."),
]
yy = top + Inches(0.75)
for t_, b_ in future:
    text(s, Inches(7.08), yy, Inches(5.45), Inches(0.75),
         [[(t_ + "  ", 11.5, True, RGBColor.from_string("FFFFFF")),
           (b_, 11, False, RGBColor.from_string("C8D9DC"))]])
    yy += Inches(0.74)
footer(s, 10)

# --- 11 · Why it matters ----------------------------------------------------------
s = blank_slide(prs); paint_bg(s, TEAL_DARK)
box(s, 0, H - Inches(0.14), W, Inches(0.14), fill=AMBER, line=None,
    shape=MSO_SHAPE.RECTANGLE)
text(s, Inches(1.0), Inches(0.7), Inches(11.3), Inches(0.4),
     [[("WHY IT MATTERS", 13, True, AMBER)]])
text(s, Inches(1.0), Inches(1.1), Inches(11.3), Inches(0.9),
     [[("Faster, more trustworthy outbreak decisions", 34, True,
        RGBColor.from_string("FFFFFF"), DISPLAY)]])
stats = [("40+", "SitReps unified", "N°001 → N°042, one queryable corpus"),
         ("5", "quality dimensions", "scored on every single report"),
         ("3", "languages", "FR · EN · PT, source-grounded"),
         ("100%", "cited answers", "SitRep N° · date · page, every time")]
cw = Inches(2.85)
for i, (n_, t_, b_) in enumerate(stats):
    x = Inches(1.0) + i * (cw + Inches(0.15))
    text(s, x, Inches(2.35), cw, Inches(0.9),
         [[(n_, 40, True, AMBER, DISPLAY)]])
    text(s, x, Inches(3.25), cw, Inches(0.8),
         [[(t_, 13.5, True, RGBColor.from_string("FFFFFF"))],
          [(b_, 10.5, False, RGBColor.from_string("9FB8BD"))]])
checks = [
    "From a pile of PDFs to a workspace analysts can interrogate.",
    "Numbers are verified, not just displayed — quality is a first-class output.",
    "Evolution is explicit: today vs. yesterday, with contradictions flagged.",
    "Extensible by design — add a report, the whole picture updates.",
]
yy = 4.55
for c_ in checks:
    text(s, Inches(1.0), Inches(yy), Inches(11.3), Inches(0.42),
         [[("✓  ", 13, True, AMBER),
           (c_, 13, False, RGBColor.from_string("D9E4E6"))]])
    yy += 0.46
text(s, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.4),
     [[("Grounded · Verified · Multilingual · Always cited", 13, True,
        RGBColor.from_string("9FB8BD"))]])

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"saved {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
