"""Generate a 5-minute overview PowerPoint for the BVD/MVE SitRep RAG application.

Run:  python scripts/build_deck.py
Output:  SitRep_Intelligence_Overview.pptx (project root)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Palette
# ----------------------------------------------------------------------------
INK = RGBColor(0x14, 0x1B, 0x2E)      # deep navy
PAPER = RGBColor(0xF7, 0xF6, 0xF2)    # warm off-white
ACCENT = RGBColor(0xD9, 0x4A, 0x38)   # terracotta (Claude-ish)
TEAL = RGBColor(0x1F, 0x6F, 0x6B)     # epidemiology green-teal
GOLD = RGBColor(0xE0, 0xA5, 0x2E)
SLATE = RGBColor(0x5A, 0x63, 0x75)    # muted slate for body
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xDD, 0xDA, 0xD0)

PRESET = [TEAL, ACCENT, GOLD, RGBColor(0x3A, 0x5A, 0x8C), RGBColor(0x7A, 0x4A, 0x8C), SLATE]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color=PAPER):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, ital) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.italic = ital
            r.font.name = "Segoe UI"
    return tb


def P(text, size, color=INK, bold=False, ital=False):
    return [(text, size, color, bold, ital)]


def kicker(s, x, y, text, color=ACCENT):
    rect(s, x, y + Emu(int(Pt(2))), Pt(22), Pt(10), color)
    txt(s, x + Pt(32), y, Inches(8), Inches(0.4),
        [P(text.upper(), 13, color, True)])


def header(s, num, title, sub=None):
    """Standard content-slide header band."""
    rect(s, 0, 0, SW, Inches(1.35), INK)
    rect(s, 0, Inches(1.35), SW, Pt(4), ACCENT)
    txt(s, Inches(0.6), Inches(0.28), Inches(1.0), Inches(0.8),
        [P(f"{num:02d}", 30, GOLD, True)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.45), Inches(0.22), Inches(11.2), Inches(0.95),
        [P(title, 27, WHITE, True)] + ([P(sub, 13, RGBColor(0xC9, 0xCE, 0xD9), False, True)] if sub else []),
        anchor=MSO_ANCHOR.MIDDLE)


def footer(s, n):
    txt(s, Inches(0.6), Inches(7.06), Inches(8), Inches(0.3),
        [P("MVE / Ebola SitRep Intelligence  ·  INSP DRC", 9, SLATE)])
    txt(s, Inches(12.0), Inches(7.06), Inches(0.9), Inches(0.3),
        [P(str(n), 9, SLATE, True)], align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, accent, title, lines, title_size=15, body_size=11.5):
    rect(s, x, y, w, h, CARD, line=LINE, line_w=1.0)
    rect(s, x, y, Pt(6), h, accent)
    txt(s, x + Inches(0.28), y + Inches(0.18), w - Inches(0.5), Inches(0.5),
        [P(title, title_size, INK, True)])
    body = [P(ln, body_size, SLATE) for ln in lines]
    txt(s, x + Inches(0.28), y + Inches(0.62), w - Inches(0.5), h - Inches(0.7),
        body, space_after=5, line_spacing=1.04)


# ----------------------------------------------------------------------------
# Slide 1 — Title
# ----------------------------------------------------------------------------
s = slide(); bg(s, INK)
rect(s, 0, 0, Inches(0.32), SH, ACCENT)
rect(s, Inches(0.32), 0, Pt(3), SH, GOLD)
kicker(s, Inches(0.95), Inches(1.15), "Outbreak situation-report intelligence", GOLD)
txt(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(2.6),
    [P("From PDF SitReps to", 46, WHITE, True),
     P("Verified Outbreak Intelligence", 46, ACCENT, True)],
    line_spacing=1.02)
txt(s, Inches(0.95), Inches(4.05), Inches(10.8), Inches(1.0),
    [P("A grounded RAG system over the 17th Ebola (MVE/BVD) outbreak reports", 17, RGBColor(0xD6,0xDA,0xE3)),
     P("from the Institut National de Santé Publique, DR Congo", 17, RGBColor(0xD6,0xDA,0xE3))],
    line_spacing=1.1)
# chips
chips = ["Ask & cite", "Verify quality", "Track evolution", "Surface gaps"]
cx = Inches(0.95)
for i, c in enumerate(chips):
    w = Inches(2.45)
    rect(s, cx, Inches(5.45), w, Inches(0.55), PRESET[i % len(PRESET)],
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, cx, Inches(5.45), w, Inches(0.55), [P(c, 13, WHITE, True)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w + Inches(0.22)
txt(s, Inches(0.95), Inches(6.55), Inches(11), Inches(0.5),
    [P("Claude  ·  Voyage embeddings  ·  Chroma  ·  FastAPI  ·  Python 3.13", 13, GOLD, True)])

# ----------------------------------------------------------------------------
# Slide 2 — The Problem
# ----------------------------------------------------------------------------
s = slide(); bg(s)
header(s, 1, "The problem", "Situation reports are the pulse of an outbreak — but they are hard to use")
intro = ("During an outbreak, INSP publishes a new PDF situation report almost every day. "
         "Analysts must read them fast, trust the numbers, and see how the picture is changing — "
         "manually, under pressure.")
txt(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.7),
    [P(intro, 14, INK)], line_spacing=1.12)

probs = [
    (ACCENT, "Data quality", ["Typos, arithmetic that doesn't add up,", "missing tables, inconsistent CFR%."]),
    (TEAL, "Verification", ["Are the figures internally consistent?", "Do today's totals agree with yesterday's?"]),
    (GOLD, "Evolution over time", ["40+ reports, N°001→N°042.", "How has the outbreak actually moved?"]),
    (RGBColor(0x3A,0x5A,0x8C), "Interaction", ["No way to just ask a question and get", "an answer with a page-level citation."]),
]
x = Inches(0.6); y = Inches(2.7); w = Inches(2.92); h = Inches(2.5); gap = Inches(0.12)
for i, (acc, t, ls) in enumerate(probs):
    card(s, x, y, w, h, acc, t, ls)
    x += w + gap
txt(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(1.2),
    [P("Consequence", 13, ACCENT, True),
     P("Decisions ride on numbers nobody has had time to validate, and trends live only in an "
       "analyst's memory of yesterday's PDF.", 14, INK)], space_after=4, line_spacing=1.12)
footer(s, 2)

# ----------------------------------------------------------------------------
# Slide 3 — The Solution / What it does
# ----------------------------------------------------------------------------
s = slide(); bg(s)
header(s, 2, "The solution", "One workspace that reads every SitRep, checks it, and answers questions")
goals = [
    (TEAL, "Ground every answer", "Responses come only from retrieved report excerpts, each cited as SitRep N° · date · page."),
    (ACCENT, "Trust the numbers", "Automatic completeness, consistency and arithmetic checks on each report — and between reports."),
    (GOLD, "See the trajectory", "Structured 'facts' datasets turn 40 PDFs into clean KPI / province / zone time-series."),
    (RGBColor(0x3A,0x5A,0x8C), "Stay current", "Drop in a new INSP link or PDF and it is parsed, extracted, indexed and ready to query."),
]
x = Inches(0.6); y = Inches(1.7); w = Inches(6.0); h = Inches(1.32)
for i, (acc, t, body) in enumerate(goals):
    col = x if i % 2 == 0 else x + w + Inches(0.18)
    row = y + (i // 2) * (h + Inches(0.2))
    card(s, col, row, w, h, acc, t, [body], body_size=12.5)
txt(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(2),
    [P("How it stays honest", 13, ACCENT, True)], space_after=6)
principles = [
    "Structured-first — exact figures come from JSON datasets, not the model's guess.",
    "Nullable by design — a missing table stays null; the system never invents a number.",
    "Tool-forced output — assessments & comparisons are returned as validated JSON, not prose.",
    "Holdout set — the 3 most recent reports are withheld from the index to test retrieval honestly.",
]
txt(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.6),
    [P("•  " + p, 13, INK) for p in principles], space_after=6, line_spacing=1.05)
footer(s, 3)

# ----------------------------------------------------------------------------
# Slide 4 — System components
# ----------------------------------------------------------------------------
s = slide(); bg(s)
header(s, 3, "System components", "Five analyst-facing capabilities over one grounded data core")
comps = [
    (ACCENT, "Grounded chat assistant", ["Multi-turn Q&A in FR / EN / PT.", "Tools: query_data (exact figures)", "+ plot_chart (trends). Always cited."]),
    (TEAL, "Outbreak summary + epi analysis", ["KPI, per-province & per-zone", "time-series: cases, deaths, CFR%,", "labs, contacts, IPC, vaccination."]),
    (GOLD, "Single-document reviewer", ["Scores 1 report on 5 dimensions,", "flags arithmetic & consistency", "issues by severity."]),
    (RGBColor(0x3A,0x5A,0x8C), "Document comparison", ["Validates two reports: totals must", "not drop, flags suspect jumps,", "contradictions, section changes."]),
    (RGBColor(0x7A,0x4A,0x8C), "Topics analysis", ["Extracts gaps & challenges per", "report, themes them, synthesises", "recurring bottlenecks."]),
    (SLATE, "SitRep library", ["Add a report via INSP link, PDF", "URL or upload → auto parse,", "extract, index, aggregate."]),
]
x0 = Inches(0.55); y0 = Inches(1.7); w = Inches(4.0); h = Inches(1.62); gx = Inches(0.13); gy = Inches(0.16)
for i, (acc, t, ls) in enumerate(comps):
    col = i % 3; row = i // 3
    card(s, x0 + col * (w + gx), y0 + row * (h + gy), w, h, acc, t, ls,
         title_size=14, body_size=11)
txt(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.4),
    [P("The data core beneath them all", 13, ACCENT, True),
     P("JSON 'facts' datasets (exact numbers)  +  Chroma vector store (semantic text search)  "
       "+  externalised prompts  +  Claude for extraction, analysis & synthesis.", 13.5, INK)],
    space_after=4, line_spacing=1.1)
footer(s, 4)

# ----------------------------------------------------------------------------
# Slide 5 — Architecture
# ----------------------------------------------------------------------------
s = slide(); bg(s)
header(s, 4, "Architecture & data flow", "Ingest once, serve many — structured facts and semantic text side by side")

def flowbox(x, y, w, h, color, title, sub, txtcol=WHITE):
    rect(s, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x + Inches(0.12), y, w - Inches(0.24), h,
        [P(title, 13, txtcol, True), P(sub, 9.5, txtcol)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)

def arrow(x, y, w):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.28))
    a.fill.solid(); a.fill.fore_color.rgb = SLATE; a.line.fill.background()
    a.shadow.inherit = False

rowy = Inches(1.85); bh = Inches(1.0)
# Source
flowbox(Inches(0.55), rowy, Inches(2.05), bh, INK, "Source SitReps", "INSP PDFs / JSON bundle")
arrow(Inches(2.68), rowy + Inches(0.36), Inches(0.5))
# Ingest
flowbox(Inches(3.25), rowy, Inches(2.2), bh, SLATE, "Ingest pipeline", "unpack · chunk · page markers")
arrow(Inches(5.53), rowy + Inches(0.36), Inches(0.5))
# Two stores
flowbox(Inches(6.1), Inches(1.5), Inches(3.0), Inches(0.85), TEAL,
        "Structured facts (JSON)", "KPI · province · zone · pillars")
flowbox(Inches(6.1), Inches(2.55), Inches(3.0), Inches(0.85), ACCENT,
        "Vector store (Chroma)", "Voyage embeddings · top-k")
arrow(Inches(9.18), rowy + Inches(0.36), Inches(0.5))
# Claude
flowbox(Inches(9.75), rowy, Inches(3.0), bh, GOLD,
        "Claude", "extract · assess · compare · synthesise", txtcol=INK)

# API + frontend band
rect(s, Inches(0.55), Inches(3.75), Inches(12.2), Inches(1.35), CARD, line=LINE)
txt(s, Inches(0.75), Inches(3.85), Inches(5), Inches(0.4),
    [P("FastAPI backend  ·  /api/*", 13, INK, True)])
eps = ["chat", "overview", "data-quality/assess", "data-quality/compare", "topics/run", "sitreps"]
ex = Inches(0.75)
for i, e in enumerate(eps):
    w = Inches(1.92)
    rect(s, ex, Inches(4.3), w, Inches(0.6), PAPER, line=LINE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, ex, Inches(4.3), w, Inches(0.6), [P("/" + e, 10.5, TEAL, True)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ex += w + Inches(0.083)
arrow(Inches(6.2), Inches(5.18), Inches(0.5))
flowbox(Inches(0.55), Inches(5.35), Inches(12.2), Inches(0.7),
        RGBColor(0x3A,0x5A,0x8C), "Claude-Design frontend",
        "6 screens · FR/EN/PT · offline sample data · Plotly charts")
txt(s, Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.6),
    [P("Caching & incrementality:", 12, ACCENT, True),
     [("per-report extraction is cached — a new upload only extracts itself, then re-aggregates the datasets.",
       12, INK, False, False)]], space_after=2, line_spacing=1.05)
footer(s, 5)

# ----------------------------------------------------------------------------
# Slide 6 — Tech stack
# ----------------------------------------------------------------------------
s = slide(); bg(s)
header(s, 5, "Tech stack", "Lean, mostly-local, API-light")
stacks = [
    (GOLD, "AI & language", ["Anthropic Claude — Sonnet 4.6", "default, Opus 4.8 for hard tasks", "Tool-forcing for structured JSON", "Prompts externalised in /prompts"]),
    (ACCENT, "Retrieval", ["Voyage AI multilingual embeddings", "Chroma local persistent vector DB", "~1400-char chunks, 200 overlap", "Top-k = 8, page-level citations"]),
    (TEAL, "Backend & data", ["FastAPI + Uvicorn (async)", "pypdf text extraction", "pandas · tiktoken · Plotly", "JSON datasets + manifest CSV"]),
    (RGBColor(0x3A,0x5A,0x8C), "Platform & quality", ["Python 3.13, uv package manager", "pytest (network-free mocks)", "ruff + black, Makefile tasks", "Claude-Design HTML frontend"]),
]
x0 = Inches(0.6); y0 = Inches(1.7); w = Inches(6.0); h = Inches(2.45); gx = Inches(0.18); gy = Inches(0.2)
for i, (acc, t, ls) in enumerate(stacks):
    col = i % 2; row = i // 2
    card(s, x0 + col * (w + gx), y0 + row * (h + gy), w, h, acc, t, ls, body_size=12.5)
footer(s, 6)

# ----------------------------------------------------------------------------
# Slide 7 — Topics analysis (deep dive)
# ----------------------------------------------------------------------------
s = slide(); bg(s)
header(s, 6, "Topics analysis", "Turning narrative 'défis & lacunes' into a strategic picture")
txt(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.75),
    [P("Each report buries its real story in prose — the gaps, constraints and unmet needs. "
       "The system extracts them per report, themes them, then synthesises what keeps recurring across the series.",
       14, INK)], line_spacing=1.12)
# two-step
rect(s, Inches(0.6), Inches(2.65), Inches(5.9), Inches(1.5), CARD, line=LINE)
rect(s, Inches(0.6), Inches(2.65), Pt(6), Inches(1.5), ACCENT)
txt(s, Inches(0.85), Inches(2.78), Inches(5.4), Inches(1.3),
    [P("1 · Extract  (gaps_extract.md)", 14, INK, True),
     P("Claude reads each SitRep and returns JSON:", 12, SLATE),
     P("{ sitrep_number, date, theme, detail }", 11.5, TEAL, True)],
    space_after=4, line_spacing=1.05)
rect(s, Inches(6.85), Inches(2.65), Inches(5.9), Inches(1.5), CARD, line=LINE)
rect(s, Inches(6.85), Inches(2.65), Pt(6), Inches(1.5), TEAL)
txt(s, Inches(7.1), Inches(2.78), Inches(5.4), Inches(1.3),
    [P("2 · Synthesise  (gaps_synthesise.md)", 14, INK, True),
     P("Aggregates all items into a narrative of the", 12, SLATE),
     P("chronic bottlenecks, cached for reuse.", 12, SLATE)],
    space_after=4, line_spacing=1.05)
txt(s, Inches(0.6), Inches(4.35), Inches(12), Inches(0.4),
    [P("THEMES TRACKED", 12, ACCENT, True)])
themes = ["Financement", "Surveillance", "Vaccination", "Sécurité", "Logistique",
          "Engagement com.", "Ressources hum.", "Laboratoire"]
tx = Inches(0.6); ty = Inches(4.75)
for i, th in enumerate(themes):
    w = Inches(2.95); col = i % 4; row = i // 4
    cx = tx + col * (w + Inches(0.12)); cy = ty + row * Inches(0.7)
    rect(s, cx, cy, w, Inches(0.58), PRESET[i % len(PRESET)],
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, cx, cy, w, Inches(0.58), [P(th, 13, WHITE, True)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.6),
    [P("Payoff:  ", 13, ACCENT, True),
     [("a coordinator sees that, e.g., financing and logistics gaps recur for weeks — "
       "evidence for where to push, not a vague impression.", 13, INK, False, False)]],
    line_spacing=1.05)
footer(s, 7)

# ----------------------------------------------------------------------------
# Slide 8 — Limitations & roadmap
# ----------------------------------------------------------------------------
s = slide(); bg(s)
header(s, 7, "Limitations & roadmap", "Honest about today's edges — and where we take it next")

def panel(x, title, acc, items):
    pw = Inches(6.0); ph = Inches(4.95); py = Inches(1.7)
    rect(s, x, py, pw, ph, CARD, line=LINE)
    rect(s, x, py, pw, Inches(0.62), acc)
    txt(s, x + Inches(0.3), py, pw - Inches(0.6), Inches(0.62),
        [P(title, 16, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    iy = py + Inches(0.85)
    for lead, body in items:
        rect(s, x + Inches(0.3), iy + Inches(0.07), Pt(9), Pt(9), acc,
             shape=MSO_SHAPE.OVAL)
        txt(s, x + Inches(0.6), iy, pw - Inches(0.9), Inches(0.95),
            [[(lead + "  ", 12.5, INK, True, False), (body, 12.5, SLATE, False, False)]],
            line_spacing=1.04)
        iy += Inches(0.79)

panel(Inches(0.6), "Current limitations", ACCENT, [
    ("API keys & cost", "needs Anthropic + Voyage keys; LLM/embedding calls are metered usage."),
    ("No OCR", "pypdf reads text only — scanned or image-only PDFs won't parse."),
    ("Format-coupled", "facts schema & prompts are tuned to the INSP / MVE report layout."),
    ("Human-in-the-loop", "LLM extraction & scoring assist analysts — they aren't authoritative."),
    ("Single-user & local", "local Chroma store, no auth, access control or multi-tenant serving."),
])
panel(Inches(6.75), "Future development", TEAL, [
    ("Multi-outbreak support", "generalise schema & prompts so any outbreak or agency can plug in."),
    ("One-click DOCX export", "share a quality assessment or comparison as Word — no copy-paste."),
    ("Deeper text analysis", "beyond topics: intervention timelines, partner engagement, response over time."),
    ("Bring-your-own key", "users supply their own API keys and pick models per deployment."),
    ("Hardening", "auth, OCR fallback, and formal retrieval benchmarking on the holdout set."),
])
footer(s, 8)

# ----------------------------------------------------------------------------
# Slide 9 — Impact / closing
# ----------------------------------------------------------------------------
s = slide(); bg(s, INK)
rect(s, 0, 0, Inches(0.32), SH, ACCENT)
kicker(s, Inches(0.9), Inches(0.7), "Why it matters", GOLD)
txt(s, Inches(0.85), Inches(1.2), Inches(11.8), Inches(1.0),
    [P("Faster, more trustworthy outbreak decisions", 36, WHITE, True)])
metrics = [
    ("40+", "SitReps unified", "N°001 → N°042, one queryable corpus"),
    ("5", "quality dimensions", "scored on every single report"),
    ("3", "languages", "FR · EN · PT, source-grounded"),
    ("100%", "cited answers", "SitRep N° · date · page, every time"),
]
x = Inches(0.85); w = Inches(2.85); gap = Inches(0.2)
for i, (big, lab, sub) in enumerate(metrics):
    rect(s, x, Inches(2.5), w, Inches(1.85), RGBColor(0x1E,0x27,0x3E),
         line=RGBColor(0x33,0x3D,0x57), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, Inches(2.5), w, Pt(5), PRESET[i % len(PRESET)])
    txt(s, x, Inches(2.68), w, Inches(0.9), [P(big, 40, GOLD, True)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.15), Inches(3.55), w - Inches(0.3), Inches(0.75),
        [P(lab, 14, WHITE, True), P(sub, 10.5, RGBColor(0xC1,0xC6,0xD2))],
        align=PP_ALIGN.CENTER, space_after=2, line_spacing=1.0)
    x += w + gap
takeaways = [
    "From a pile of PDFs to a workspace analysts can interrogate.",
    "Numbers are verified, not just displayed — quality is a first-class output.",
    "Evolution is explicit: today vs. yesterday, with contradictions flagged.",
    "Extensible by design — add a report, the whole picture updates.",
]
txt(s, Inches(0.85), Inches(4.7), Inches(11.6), Inches(1.6),
    [P("✓  " + t, 14, RGBColor(0xE9,0xEC,0xF2)) for t in takeaways],
    space_after=7, line_spacing=1.05)
rect(s, Inches(0.85), Inches(6.55), Inches(11.6), Pt(3), ACCENT)
txt(s, Inches(0.85), Inches(6.7), Inches(11.6), Inches(0.5),
    [P("Grounded · Verified · Multilingual · Always cited", 15, GOLD, True)])

out = "SitRep_Intelligence_Overview.pptx"
try:
    prs.save(out)
except PermissionError:
    out = "SitRep_Intelligence_Overview_v2.pptx"
    prs.save(out)
    print("(original file was open/locked — saved to a new file instead)")
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
